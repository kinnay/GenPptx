
from PIL import Image, ImageOps

from dataclasses import dataclass
from pathlib import Path

from tkinter.messagebox import showwarning, showerror

import pptx
import pptx.presentation
import pptx.util
import sys

import logging
logger = logging.getLogger("GenPptx")

logging.basicConfig(level=logging.INFO, format="%(message)s")


VERSION = "1.0.0"


@dataclass
class SongReference:
    id: str
    verses: list[int] | None


@dataclass
class SongImage:
    image: Image.Image
    path: Path

    def cropped(self) -> Image.Image:
        brightness = self.image.split()[0]
        darkness = ImageOps.invert(brightness)
        mask = darkness.getbbox()
        return self.image.crop(mask)


class ImageDB:
    _songs: dict[str, dict[int, dict[int, SongImage]]]

    def __init__(self):
        self._songs = {}
    
    def load(self, path: Path) -> None:
        try:
            image = Image.open(path)
        except Exception:
            logger.exception(f"Fout bij het openen van {path}.")
            showwarning(
                "Fout in afbeelding",
                f"Er is een fout opgetreden bij het openen van {path}."
            )
            return
        
        parts = path.stem.split("-")

        id = parts[3]
        if len(parts) == 6:
            verse = 1
            slide = int(parts[5])
        else:
            verse = int(parts[6])
            slide = int(parts[7])
        
        if id not in self._songs:
            self._songs[id] = {}
        
        if verse not in self._songs[id]:
            self._songs[id][verse] = {}
        
        self._songs[id][verse][slide] = SongImage(image, path)
    
    def find(self, reference: SongReference) -> list[SongImage]:
        song_images = self._songs.get(reference.id)
        if not song_images:
            logger.warning(
                f"Kon geen afbeeldingen vinden voor lied {reference.id}."
            )
            return []
        
        verses = reference.verses
        if verses is None:
            verses = sorted(song_images)
        
        images = []
        for verse in verses:
            if verse not in song_images:
                logger.warning(
                    f"Kon geen afbeeldingen vinden voor couplet {verse} van "
                    f"lied {reference.id}."
                )
                continue

            for slide, image in sorted(song_images[verse].items()):
                images.append(image)
        
        return images
    
    def find_all(self, references: list[SongReference]) -> list[SongImage]:
        images = []
        for reference in references:
            images += self.find(reference)
        return images


class Powerpoint:
    _presentation: pptx.presentation.Presentation

    def __init__(self):
        self._presentation = pptx.Presentation()
        self._presentation.slide_width = pptx.util.Cm(16)
        self._presentation.slide_height = pptx.util.Cm(9)
    
    def add_image(self, image: Image.Image, path: Path) -> None:
        ratio = image.height / image.width

        max_width = pptx.util.Cm(15)
        max_height = pptx.util.Cm(8)

        target_width = max_width
        target_height = max_width * ratio
        if target_height > max_height:
            target_width *= max_height / target_height
            target_height = max_height
        
        target_x = (self._presentation.slide_width - target_width) / 2
        target_y = (self._presentation.slide_height - target_height) / 2

        blank_layout = self._presentation.slide_layouts[6]

        slide = self._presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(path), target_x, target_y, target_width, target_height
        )
    
    def save(self, path: Path):
        self._presentation.save(str(path))


def check_image_name(name: str) -> bool:
    # Possible filenames:
    # * nl-lb-projectie-<nummer>-muziek-<dia>
    # * nl-lb-projectie-<nummer>-muziek-couplet-<couplet>-<dia>
    
    if not name.startswith("nl-lb-projectie-"):
        return False

    parts = name.split("-")
    if len(parts) == 6:
        if parts[4] != "muziek": return False
        if not parts[5].isdecimal(): return False
    elif len(parts) == 8:
        if parts[4] != "muziek": return False
        if parts[5] != "couplet": return False
        if not parts[6].isdecimal(): return False
        if not parts[7].isdecimal(): return False
    else:
        return False
    
    return True

def load_song_references(path: Path) -> list[SongReference]:
    with open(path) as f:
        lines = [line.strip() for line in f if line.strip()]
    
    references = []
    for line in lines:
        parts = line.split(maxsplit=1)

        id = parts[0]
        verses = None
        if len(parts) > 1:
            verse_parts = [part.strip() for part in parts[1].split(",")]
            if not all(part.isdecimal() for part in verse_parts):
                logger.warning(f"Regel {line} is ongeldig.")
                showwarning(
                    "Ongeldige regel",
                    f'De regel "{line}" in Liederen.txt heeft niet het juiste '
                        "formaat en wordt daarom genegeerd."
                )
                continue

            verses = [int(part) for part in verse_parts]
        
        references.append(SongReference(id, verses))
    return references

def create_powerpoint(
    database: ImageDB, references: list[SongReference], base_path: Path
) -> None:
    powerpoint = Powerpoint()

    images = database.find_all(references)
    logger.info(f"{len(images)} afbeeldingen worden verwerkt.")

    crops_folder = base_path / "crops"
    crops_folder.mkdir(exist_ok=True)

    for image in images:
        crop_path = crops_folder / image.path.name

        cropped = image.cropped()
        cropped.save(crop_path)

        powerpoint.add_image(cropped, crop_path)
    
    powerpoint_path = base_path / "Liederen.pptx"
    powerpoint.save(powerpoint_path)

    logger.info(f"De PowerPoint staat klaar in {powerpoint_path}!")


def main() -> None:
    print(f"GenPptx versie {VERSION}")

    input_path = Path("Liederen")
    if len(sys.argv) >= 2:
        input_path = Path(sys.argv[1])
    
    parent_path = input_path.parent
    
    if not input_path.exists():
        logger.warning(f"De map {input_path} kon niet worden gevonden.")
        showwarning(
            "Geen liederen gevonden",
            f'De map "{input_path}" kon niet worden gevonden.',
            detail = \
                'Zorg ervoor dat er een map "Liederen" in de map van dit ' \
                "programma staat, of sleep een andere map van het systeem naar "
                f"{Path(sys.argv[0]).name} toe."
        )
        return
    
    logger.info(f"De map {input_path} wordt naar liederen doorzocht.")

    database = ImageDB()
    for image_path in input_path.glob("nl-lb-projectie-*.png"):
        if not check_image_name(image_path.stem):
            logger.warning(f"Het bestand {image_path.name} wordt genegeerd.")
            continue

        database.load(image_path)

    reference_path = parent_path / "Liederen.txt"
    if not reference_path.exists():
        logger.warning(f"{reference_path} kon niet worden gevonden.")
        showwarning(
            "Liederen.txt kon niet worden gevonden. Zorg ervoor dat dit "
            f"bestand zich in de map {parent_path} bevindt."
        )
        return
    
    references = load_song_references(reference_path)
    logger.info(f"Er zijn {len(references)} regels gevonden in Liederen.txt.")

    create_powerpoint(database, references, parent_path)


if __name__ == "__main__":
    try:
        main()
    except Exception:
        logger.exception("Een onverwachte fout is opgetreden")
        showerror(
            "Onverwachte fout",
            "Een onverwachte fout is opgetreden."
        )
    
    input("Druk op enter om dit venster te sluiten...")
